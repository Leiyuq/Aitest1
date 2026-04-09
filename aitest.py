"""
测试用例智构系统 - 轻量级 RAG（NumPy TF‑IDF）  无用户登录、项目知识库
阿里云OSS持久化存储
"""
import streamlit as st
st.set_page_config(page_title="测试用例智构系统", layout="wide")

import os
import re
import json
import pickle
import io
import time
import warnings
import tempfile
from pathlib import Path
from datetime import datetime
from typing import List, Dict, Optional, Callable
from dataclasses import dataclass
from concurrent.futures import ThreadPoolExecutor, TimeoutError as FuturesTimeoutError

import numpy as np
import pandas as pd
from openai import OpenAI
import oss2
from oss2.exceptions import NoSuchKey

# 可选依赖：用于解析特殊格式
HAS_DOCX = False
try:
    import docx
    HAS_DOCX = True
except ImportError:
    pass

try:
    import PyPDF2
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    from openpyxl import load_workbook
    HAS_EXCEL = True
except ImportError:
    HAS_EXCEL = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import xmindparser
    HAS_XMIND = True
except ImportError:
    HAS_XMIND = False

try:
    import jieba
    USE_JIEBA = True
except ImportError:
    USE_JIEBA = False
    jieba = None

warnings.filterwarnings('ignore')


# ====================== 配置 ======================
@dataclass
class ModelConfig:
    name: str
    api_key: str
    base_url: str
    model: str

class AppConfig:
    PAGE_TITLE = "测试用例智构系统"
    LAYOUT = "wide"
    BASE_ROOT = "knowledge_projects"        # OSS 中的根目录前缀
    ALLOWED_FILE_TYPES = {
        "txt", "csv", "docx", "doc", "pdf", "xlsx", "xls", "pptx", "xmind", "md"
    }
    RAG_TOP_K = 5
    RAG_SIMILARITY_THRESHOLD = 0.2
    API_TIMEOUT = 60
    API_MAX_RETRIES = 2
    API_RETRY_DELAY = 2

    MODELS = {
        "qwen": ModelConfig(
            name="通义千问",
            api_key="sk-66b40867237b4e589c81c0255ff94b36",
            base_url="https://dashscope.aliyuncs.com/compatible-mode/v1",
            model="qwen-plus"
        ),
        "local": ModelConfig(
            name="本地离线",
            api_key="none",
            base_url="none",
            model="local"
        )
    }

    @classmethod
    def get_model_list(cls):
        return [f"{m.name}({k})" for k, m in cls.MODELS.items() if m.api_key or k == "local"]


# ====================== 阿里云 OSS 辅助函数 ======================
def get_oss_bucket():
    """从 st.secrets 获取 OSS 配置并返回 Bucket 对象"""
    try:
        auth = oss2.Auth(
            st.secrets["OSS_ACCESS_KEY_ID"],
            st.secrets["OSS_ACCESS_KEY_SECRET"]
        )
        bucket = oss2.Bucket(
            auth,
            st.secrets["OSS_ENDPOINT"],
            st.secrets["OSS_BUCKET_NAME"]
        )
        return bucket
    except Exception as e:
        st.error(f"OSS 配置错误: {e}")
        st.stop()


# ====================== 轻量级 TF‑IDF 向量检索（支持 OSS） ======================
class SimpleVectorStore:
    def __init__(self, store_dir: str, oss_bucket, oss_prefix: str):
        """
        store_dir: 本地临时目录，用于存放下载的索引文件
        oss_bucket: OSS Bucket 对象
        oss_prefix: OSS 中该项目的向量存储前缀（如 "knowledge_projects/项目名/vector_store"）
        """
        self.store_dir = store_dir
        self.oss_bucket = oss_bucket
        self.oss_prefix = oss_prefix
        Path(store_dir).mkdir(parents=True, exist_ok=True)

        self.index_file = os.path.join(store_dir, "index.json")
        self.vectors_file = os.path.join(store_dir, "vectors.npy")
        self.metadata_file = os.path.join(store_dir, "metadata.pkl")
        self.documents = []
        self.vocab = {}
        self.idf = {}
        self.vectors = None

    def _oss_key(self, filename: str) -> str:
        """生成 OSS 对象键"""
        return f"{self.oss_prefix}/{filename}"

    def _download_from_oss(self, local_path: str, object_key: str) -> bool:
        """从 OSS 下载文件到本地，如果不存在返回 False"""
        try:
            self.oss_bucket.get_object_to_file(object_key, local_path)
            return True
        except NoSuchKey:
            return False

    def _upload_to_oss(self, local_path: str, object_key: str):
        """上传本地文件到 OSS"""
        self.oss_bucket.put_object_from_file(object_key, local_path)

    def _delete_oss_object(self, object_key: str):
        """删除 OSS 对象"""
        try:
            self.oss_bucket.delete_object(object_key)
        except Exception:
            pass

    @staticmethod
    def _tokenize(text: str) -> List[str]:
        text = text.lower()
        if USE_JIEBA and jieba:
            words = jieba.lcut(text)
            return [w for w in words if len(w.strip()) > 1 or w.isalnum()]
        else:
            return re.findall(r'[\u4e00-\u9fa5]{2,}|[a-z]{2,}', text)

    @staticmethod
    def _clean_query(query: str) -> str:
        query = re.sub(r'\b[A-Za-z0-9]+[_-][A-Za-z0-9]+\b', '', query)
        query = re.sub(r'\b[A-Z]{2,}[0-9]+\b', '', query)
        query = re.sub(r'\b[A-Z]+-\d+\b', '', query)
        query = re.sub(r'\b\d{4}[-/]\d{1,2}[-/]\d{1,2}\b', '', query)
        query = re.sub(r'\b\d{1,2}[-/]\d{1,2}[-/]\d{4}\b', '', query)
        query = re.sub(r'\b\d{8}\b', '', query)
        query = re.sub(r'\s+', ' ', query).strip()
        return query

    def _build_vocab_and_idf(self, all_texts: List[str]):
        doc_count = {}
        for text in all_texts:
            unique_words = set(self._tokenize(text))
            for w in unique_words:
                doc_count[w] = doc_count.get(w, 0) + 1

        total_docs = len(all_texts)
        self.vocab = {}
        self.idf = {}
        for w, df in doc_count.items():
            if 0.01 < df / total_docs < 0.8:
                idx = len(self.vocab)
                self.vocab[w] = idx
                self.idf[w] = np.log((total_docs + 1) / (df + 0.5))

    def _text_to_vector(self, text: str) -> np.ndarray:
        words = self._tokenize(text)
        if not words or not self.vocab:
            return np.zeros(len(self.vocab))
        tf = {}
        for w in words:
            if w in self.vocab:
                tf[w] = tf.get(w, 0) + 1
        vec = np.zeros(len(self.vocab))
        for w, cnt in tf.items():
            idx = self.vocab[w]
            tf_val = np.log1p(cnt / len(words))
            vec[idx] = tf_val * self.idf.get(w, 1.0)
        norm = np.linalg.norm(vec)
        if norm > 0:
            vec /= norm
        return vec

    def build_index(self, documents: List[Dict], files_info: List[Dict] = None) -> bool:
        try:
            if not documents:
                return False
            self.documents = [{"id": i, "text": d["content"], "metadata": d.get("metadata", {})}
                              for i, d in enumerate(documents)]
            texts = [d["text"] for d in self.documents]
            self._build_vocab_and_idf(texts)
            vectors = [self._text_to_vector(t) for t in texts]
            self.vectors = np.array(vectors)

            # 保存到本地临时文件
            with open(self.index_file, 'w', encoding='utf-8') as f:
                json.dump({
                    "documents": self.documents,
                    "vocab": self.vocab,
                    "idf": self.idf
                }, f, ensure_ascii=False)
            np.save(self.vectors_file, self.vectors)
            if files_info:
                with open(self.metadata_file, 'wb') as f:
                    pickle.dump({"files_info": files_info, "updated": datetime.now().isoformat()}, f)

            # 上传到 OSS
            self._upload_to_oss(self.index_file, self._oss_key("index.json"))
            self._upload_to_oss(self.vectors_file, self._oss_key("vectors.npy"))
            if files_info and os.path.exists(self.metadata_file):
                self._upload_to_oss(self.metadata_file, self._oss_key("metadata.pkl"))
            return True
        except Exception as e:
            st.error(f"构建索引失败: {e}")
            return False

    def load_index(self) -> bool:
        try:
            if not self._download_from_oss(self.index_file, self._oss_key("index.json")):
                return False
            if not self._download_from_oss(self.vectors_file, self._oss_key("vectors.npy")):
                return False
            self._download_from_oss(self.metadata_file, self._oss_key("metadata.pkl"))

            with open(self.index_file, 'r', encoding='utf-8') as f:
                data = json.load(f)
                self.documents = data["documents"]
                self.vocab = data["vocab"]
                self.idf = data["idf"]
            self.vectors = np.load(self.vectors_file)
            return True
        except Exception as e:
            st.warning(f"加载索引失败: {e}")
            return False

    def search(self, query: str, top_k: int = 5) -> List[Dict]:
        if self.vectors is None or len(self.documents) == 0:
            return []
        cleaned_query = self._clean_query(query)
        if not cleaned_query:
            cleaned_query = query
        q_vec = self._text_to_vector(cleaned_query)
        scores = np.dot(self.vectors, q_vec)
        top_idx = np.argsort(scores)[-top_k:][::-1]
        results = []
        for idx in top_idx:
            sim = float(scores[idx])
            if sim >= AppConfig.RAG_SIMILARITY_THRESHOLD:
                doc = self.documents[idx]
                results.append({
                    "content": doc["text"],
                    "metadata": doc["metadata"],
                    "similarity": sim
                })
        return results

    def is_index_valid(self, current_files_info: List[Dict]) -> bool:
        """检查 OSS 上的元数据文件是否与当前文件列表一致"""
        temp_meta = os.path.join(self.store_dir, "metadata_check.pkl")
        if not self._download_from_oss(temp_meta, self._oss_key("metadata.pkl")):
            return False
        try:
            with open(temp_meta, 'rb') as f:
                meta = pickle.load(f)
            saved = {f["name"]: f for f in meta["files_info"]}
            curr = {f["name"]: f for f in current_files_info}
            if set(saved.keys()) != set(curr.keys()):
                return False
            for name, cf in curr.items():
                sf = saved[name]
                if sf.get("size_bytes") != cf.get("size_bytes") or sf.get("mtime") != cf.get("mtime"):
                    return False
            return True
        finally:
            if os.path.exists(temp_meta):
                os.remove(temp_meta)

    def get_built_files(self) -> set:
        if not self.documents:
            return set()
        return {doc["metadata"].get("source") for doc in self.documents if "source" in doc["metadata"]}

    def get_stats(self) -> Dict:
        return {
            "documents": len(self.documents),
            "vocab_size": len(self.vocab),
            "has_index": self.vectors is not None
        }

    def remove_file(self, filename: str) -> bool:
        """从当前内存索引中删除指定文件的分块，并重建索引，然后上传到 OSS"""
        if not self.documents:
            return False
        new_docs = [doc for doc in self.documents if doc["metadata"].get("source") != filename]
        if len(new_docs) == len(self.documents):
            return False
        if not new_docs:
            self.documents = []
            self.vocab = {}
            self.idf = {}
            self.vectors = None
            for fname in ["index.json", "vectors.npy", "metadata.pkl"]:
                self._delete_oss_object(self._oss_key(fname))
            for f in [self.index_file, self.vectors_file, self.metadata_file]:
                if os.path.exists(f):
                    os.remove(f)
            return True
        # 重建索引
        texts = [doc["text"] for doc in new_docs]
        self._build_vocab_and_idf(texts)
        vectors = [self._text_to_vector(t) for t in texts]
        self.vectors = np.array(vectors)
        self.documents = new_docs
        with open(self.index_file, 'w', encoding='utf-8') as f:
            json.dump({
                "documents": self.documents,
                "vocab": self.vocab,
                "idf": self.idf
            }, f, ensure_ascii=False)
        np.save(self.vectors_file, self.vectors)
        self._upload_to_oss(self.index_file, self._oss_key("index.json"))
        self._upload_to_oss(self.vectors_file, self._oss_key("vectors.npy"))
        return True

    def clear_index(self):
        self.documents = []
        self.vocab = {}
        self.idf = {}
        self.vectors = None
        for fname in ["index.json", "vectors.npy", "metadata.pkl"]:
            self._delete_oss_object(self._oss_key(fname))
        for f in [self.index_file, self.vectors_file, self.metadata_file]:
            if os.path.exists(f):
                os.remove(f)

    def update_metadata(self, files_info: List[Dict]):
        """单独更新 metadata.pkl 到 OSS"""
        if not files_info:
            self._delete_oss_object(self._oss_key("metadata.pkl"))
            return
        with open(self.metadata_file, 'wb') as f:
            pickle.dump({"files_info": files_info, "updated": datetime.now().isoformat()}, f)
        self._upload_to_oss(self.metadata_file, self._oss_key("metadata.pkl"))


# ====================== 项目管理（基于 OSS 前缀） ======================
class ProjectManager:
    @staticmethod
    def get_oss_bucket():
        return get_oss_bucket()

    @staticmethod
    def get_all_projects() -> List[str]:
        """通过列出 OSS 前缀获取所有项目名称"""
        bucket = ProjectManager.get_oss_bucket()
        prefix = f"{AppConfig.BASE_ROOT}/"
        try:
            projects = []
            for obj in oss2.ObjectIterator(bucket, prefix=prefix, delimiter='/'):
                # obj.key 格式如 "knowledge_projects/项目名/"
                if obj.key.endswith('/'):
                    proj_name = obj.key[len(prefix):-1]
                    if proj_name:
                        projects.append(proj_name)
            return sorted(projects)
        except Exception:
            return []

    @staticmethod
    def create_project(project_name: str) -> bool:
        if not project_name or not project_name.strip():
            return False
        project_name = project_name.strip()
        if not re.match(r'^[\u4e00-\u9fa5a-zA-Z0-9_-]+$', project_name):
            st.error("项目名称只能包含中文、字母、数字、下划线和中划线")
            return False
        existing = ProjectManager.get_all_projects()
        if project_name in existing:
            return False
        bucket = ProjectManager.get_oss_bucket()
        marker_key = f"{AppConfig.BASE_ROOT}/{project_name}/.project"
        try:
            bucket.put_object(marker_key, b'')
            bucket.put_object(f"{AppConfig.BASE_ROOT}/{project_name}/knowledge_base/.keep", b'')
            bucket.put_object(f"{AppConfig.BASE_ROOT}/{project_name}/vector_store/.keep", b'')
            return True
        except Exception:
            return False

    @staticmethod
    def get_project_path(project_name: str) -> str:
        return f"{AppConfig.BASE_ROOT}/{project_name}"

    @staticmethod
    def get_kb_path(project_name: str) -> str:
        return f"{AppConfig.BASE_ROOT}/{project_name}/knowledge_base"

    @staticmethod
    def get_vector_store_path(project_name: str) -> str:
        return f"{AppConfig.BASE_ROOT}/{project_name}/vector_store"


# ====================== 知识库封装（OSS 存储） ======================
class EnhancedKnowledgeBase:
    def __init__(self, project_name: str):
        self.project_name = project_name
        self.oss_bucket = get_oss_bucket()
        self.kb_prefix = ProjectManager.get_kb_path(project_name)      # OSS 前缀
        self.vector_prefix = ProjectManager.get_vector_store_path(project_name)

        # 临时目录
        self.temp_dir = Path(tempfile.gettempdir()) / f"kb_{project_name}"
        self.temp_dir.mkdir(parents=True, exist_ok=True)

        self.vector_store = SimpleVectorStore(
            str(self.temp_dir / "vector_store"),
            self.oss_bucket,
            self.vector_prefix
        )
        self.chunk_size = 800
        self.chunk_overlap = 200
        self.index_loaded = False
        self._auto_load()

    def _auto_load(self):
        files = self.get_file_list(with_metadata=True)
        if files and self.vector_store.is_index_valid(files):
            if self.vector_store.load_index():
                self.index_loaded = True

    def refresh_index(self) -> Dict:
        self.index_loaded = False
        files = self.get_file_list(with_metadata=True)
        if not files:
            self.vector_store.clear_index()
            return {"status": "error", "message": "知识库无文件"}
        if self.vector_store.load_index():
            self.index_loaded = True
            stats = self.vector_store.get_stats()
            built_files = self.vector_store.get_built_files()
            all_files = {f["name"] for f in files}
            pending = all_files - built_files
            if pending:
                msg = f"已加载索引，包含 {stats['documents']} 个片段。注意：有 {len(pending)} 个文件未构建，请点击「构建知识库」增量添加。"
            else:
                msg = f"已加载索引，包含 {stats['documents']} 个知识片段"
            return {"status": "success", "message": msg}
        return {"status": "error", "message": "未找到有效索引，请先构建知识库"}

    def get_file_list(self, with_metadata=False):
        files = []
        prefix = f"{self.kb_prefix}/"
        try:
            for obj in oss2.ObjectIterator(self.oss_bucket, prefix=prefix):
                key = obj.key
                filename = key[len(prefix):]
                if not filename or filename.startswith('.') or filename.endswith('/'):
                    continue
                ext = filename.split('.')[-1].lower()
                if ext in AppConfig.ALLOWED_FILE_TYPES:
                    size_kb = obj.size / 1024
                    info = {"name": filename, "size": f"{size_kb:.1f}KB"}
                    if with_metadata:
                        # 兼容 last_modified 可能是 datetime 或 时间戳
                        if hasattr(obj.last_modified, 'timestamp'):
                            mtime = obj.last_modified.timestamp()
                        else:
                            mtime = float(obj.last_modified)
                        info.update({"size_bytes": obj.size, "mtime": mtime})
                    files.append(info)
            return files
        except Exception as e:
            st.error(f"OSS 列表失败: {e}")
            return []

    def upload_file(self, filename: str, content: bytes) -> bool:
        key = f"{self.kb_prefix}/{filename}"
        try:
            # 检查是否存在
            self.oss_bucket.head_object(key)
            return False
        except NoSuchKey:
            pass
        try:
            self.oss_bucket.put_object(key, content)
            self.index_loaded = False
            return True
        except Exception:
            return False

    def delete_file(self, filename: str):
        key = f"{self.kb_prefix}/{filename}"
        try:
            self.oss_bucket.delete_object(key)
        except Exception:
            pass
        self.vector_store.remove_file(filename)
        remaining_files = self.get_file_list(with_metadata=True)
        self.vector_store.update_metadata(remaining_files)
        self.index_loaded = len(remaining_files) > 0 and self.vector_store.load_index()

    # ---------- 文档解析（从 OSS 读取字节流） ----------
    def _read_file_from_oss(self, filename: str) -> bytes:
        key = f"{self.kb_prefix}/{filename}"
        return self.oss_bucket.get_object(key).read()

    def _extract_text_from_bytes(self, filename: str, content_bytes: bytes) -> str:
        ext = filename.split('.')[-1].lower()
        # 对于需要文件路径的解析器，写入临时文件
        if ext in ['docx', 'doc', 'pdf', 'xlsx', 'xls', 'pptx', 'xmind']:
            with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as tmp:
                tmp.write(content_bytes)
                tmp_path = tmp.name
            try:
                return self._extract_text_from_path(tmp_path, ext)
            finally:
                os.unlink(tmp_path)
        else:
            # txt, csv, json, md 直接解码
            if ext in ['txt', 'md', 'json']:
                return content_bytes.decode('utf-8', errors='ignore')
            elif ext == 'csv':
                df = pd.read_csv(io.BytesIO(content_bytes))
                return df.to_string()
            else:
                return f"暂无法识别该类型 (.{ext})，请联系管理员"

    def _extract_text_from_path(self, filepath: str, ext: str) -> str:
        parsers = {
            "docx": lambda p: '\n'.join([para.text for para in docx.Document(p).paragraphs]) if HAS_DOCX else "请安装 python-docx",
            "doc": lambda p: '\n'.join([para.text for para in docx.Document(p).paragraphs]) if HAS_DOCX else "请安装 python-docx",
            "pdf": lambda p: '\n'.join([page.extract_text() for page in PyPDF2.PdfReader(p).pages]) if HAS_PDF else "请安装 PyPDF2",
            "xlsx": lambda p: self._parse_excel(p),
            "xls": lambda p: self._parse_excel(p),
            "pptx": lambda p: '\n'.join([shape.text for slide in Presentation(p).slides for shape in slide.shapes if hasattr(shape, "text")]) if HAS_PPTX else "请安装 python-pptx",
            "xmind": lambda p: self._parse_xmind(p),
        }
        if ext in parsers:
            try:
                return parsers[ext](filepath)
            except Exception as e:
                return f"提取失败: {e}"
        return f"暂无法识别该类型 (.{ext})，请联系管理员"

    @staticmethod
    def _parse_excel(filepath: str) -> str:
        if not HAS_EXCEL:
            return "请安装 openpyxl"
        wb = load_workbook(filepath, data_only=True)
        all_text = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = ' '.join([str(cell) for cell in row if cell is not None])
                if row_text:
                    all_text.append(row_text)
        return '\n'.join(all_text)

    @staticmethod
    def _parse_xmind(filepath: str) -> str:
        if not HAS_XMIND:
            return "请安装 xmindparser"
        data = xmindparser.xmind_to_dict(filepath)
        def extract(node):
            texts = []
            if isinstance(node, dict):
                if 'title' in node:
                    texts.append(node['title'])
                if 'note' in node and node['note']:
                    texts.append(node['note'])
                if 'topics' in node:
                    for sub in node['topics']:
                        texts.extend(extract(sub))
            elif isinstance(node, list):
                for item in node:
                    texts.extend(extract(item))
            return texts
        return '\n'.join(extract(data))

    # ---------- 分块与构建 ----------
    def _chunk_text(self, text: str, source: str) -> List[Dict]:
        prefixed_text = f"[文件：{source}]\n{text}"
        chunks = []
        paragraphs = prefixed_text.split('\n')
        current = ""
        chunk_id = 0
        for para in paragraphs:
            para = para.strip()
            if not para:
                continue
            if len(para) > self.chunk_size:
                sentences = re.split(r'[。！？；]', para)
                for sentence in sentences:
                    sentence = sentence.strip()
                    if not sentence:
                        continue
                    if len(current) + len(sentence) <= self.chunk_size:
                        current += sentence + "。"
                    else:
                        if current:
                            chunks.append({"content": current.strip(), "metadata": {"source": source, "chunk_id": chunk_id}})
                            chunk_id += 1
                        overlap = current[-self.chunk_overlap:] if current else ""
                        current = overlap + sentence + "。"
            else:
                if len(current) + len(para) <= self.chunk_size:
                    current += para + "\n"
                else:
                    if current:
                        chunks.append({"content": current.strip(), "metadata": {"source": source, "chunk_id": chunk_id}})
                        chunk_id += 1
                    overlap = current[-self.chunk_overlap:] if current else ""
                    current = overlap + para + "\n"
        if current:
            chunks.append({"content": current.strip(), "metadata": {"source": source, "chunk_id": chunk_id}})
        return chunks

    def build_knowledge_base(self, progress_callback: Optional[Callable] = None, force=False) -> Dict:
        try:
            files = self.get_file_list(with_metadata=True)
            if not files:
                return {"status": "error", "message": "无文件"}

            if self.index_loaded:
                built_files = self.vector_store.get_built_files()
            else:
                if self.vector_store.load_index():
                    self.index_loaded = True
                    built_files = self.vector_store.get_built_files()
                else:
                    built_files = set()

            files_to_build = [f for f in files if f["name"] not in built_files]

            if not force and not files_to_build:
                stats = self.vector_store.get_stats()
                return {"status": "success", "message": f"所有文档已构建（共 {stats['documents']} 个知识片段）", "chunks": stats['documents']}

            # 增量构建
            if built_files and not force:
                new_chunks = []
                total_to_build = len(files_to_build)
                for i, f in enumerate(files_to_build):
                    if progress_callback:
                        progress_callback(i/total_to_build, f"处理 {f['name']} (解析中...)")
                    content_bytes = self._read_file_from_oss(f["name"])
                    text = self._extract_text_from_bytes(f["name"], content_bytes)
                    if text.startswith("请安装") or text.startswith("暂无法识别") or text.startswith("提取失败"):
                        continue
                    chunks = self._chunk_text(text, f["name"])
                    new_chunks.extend(chunks)
                    if progress_callback:
                        progress_callback(i/total_to_build, f"处理 {f['name']} → 生成 {len(chunks)} 个片段")

                if not new_chunks:
                    return {"status": "error", "message": "没有可构建的有效内容，请检查文件格式"}

                all_docs = self.vector_store.documents + [
                    {"id": i+len(self.vector_store.documents), "text": c["content"], "metadata": c["metadata"]}
                    for i, c in enumerate(new_chunks)
                ]
                texts = [doc["text"] for doc in all_docs]
                self.vector_store._build_vocab_and_idf(texts)
                vectors = [self.vector_store._text_to_vector(t) for t in texts]
                self.vector_store.vectors = np.array(vectors)
                self.vector_store.documents = all_docs
                # 保存到本地并上传 OSS
                with open(self.vector_store.index_file, 'w', encoding='utf-8') as f:
                    json.dump({
                        "documents": self.vector_store.documents,
                        "vocab": self.vector_store.vocab,
                        "idf": self.vector_store.idf
                    }, f, ensure_ascii=False)
                np.save(self.vector_store.vectors_file, self.vector_store.vectors)
                self.vector_store._upload_to_oss(self.vector_store.index_file, self.vector_store._oss_key("index.json"))
                self.vector_store._upload_to_oss(self.vector_store.vectors_file, self.vector_store._oss_key("vectors.npy"))
                files_info = [{"name": f["name"], "size_bytes": f["size_bytes"], "mtime": f["mtime"]} for f in files]
                self.vector_store.update_metadata(files_info)
                self.index_loaded = True
                return {"status": "success", "message": f"增量构建成功，新增 {len(files_to_build)} 个文件，{len(new_chunks)} 个知识片段", "files": len(files), "chunks": len(all_docs)}
            else:
                # 全量构建
                all_chunks = []
                missing_deps = set()
                unsupported_files = []
                total_files = len(files)
                for i, f in enumerate(files):
                    if progress_callback:
                        progress_callback(i/total_files, f"处理 {f['name']} (解析中...)")
                    content_bytes = self._read_file_from_oss(f["name"])
                    text = self._extract_text_from_bytes(f["name"], content_bytes)
                    if text.startswith("请安装"):
                        missing_deps.add(text)
                        continue
                    if text.startswith("暂无法识别该类型"):
                        unsupported_files.append(f"{f['name']}: {text}")
                        continue
                    if text.startswith("提取失败"):
                        unsupported_files.append(f"{f['name']}: {text}")
                        continue
                    chunks = self._chunk_text(text, f["name"])
                    all_chunks.extend(chunks)
                    if progress_callback:
                        progress_callback(i/total_files, f"处理 {f['name']} → 生成 {len(chunks)} 个片段")
                if missing_deps:
                    st.warning("\n".join(missing_deps))
                if unsupported_files:
                    for err in unsupported_files:
                        st.error(err)
                    if not all_chunks:
                        return {"status": "error", "message": "所有文件均无法解析，请检查依赖或文件格式"}
                if progress_callback:
                    progress_callback(0.95, "构建向量索引...")
                files_info = [{"name": f["name"], "size_bytes": f["size_bytes"], "mtime": f["mtime"]} for f in files]
                ok = self.vector_store.build_index(all_chunks, files_info)
                if ok:
                    self.index_loaded = True
                    return {"status": "success", "message": f"构建成功，共 {len(files)} 个文件，{len(all_chunks)} 个知识片段", "files": len(files), "chunks": len(all_chunks)}
                else:
                    return {"status": "error", "message": "索引失败"}
        except Exception as e:
            st.error(f"构建知识库时发生未预期的错误: {str(e)}")
            return {"status": "error", "message": str(e)}

    def search_knowledge(self, query: str, top_k=5):
        if not self.index_loaded and not self.vector_store.load_index():
            return []
        return self.vector_store.search(query, top_k)

    def get_knowledge_context(self, query: str, max_chunks=5) -> str:
        results = self.search_knowledge(query, max_chunks)
        if not results:
            return ""
        parts = []
        for i, r in enumerate(results, 1):
            parts.append(f"【知识片段 {i}】来源：{r['metadata'].get('source','未知')} 相关度：{r['similarity']:.2%}\n{r['content']}")
        return "\n\n---\n\n".join(parts)

    def get_built_files_safe(self) -> set:
        if self.index_loaded:
            return self.vector_store.get_built_files()
        else:
            if self.vector_store.load_index():
                self.index_loaded = True
                return self.vector_store.get_built_files()
            return set()


# ====================== LLM 服务（与原来完全相同，省略重复代码） ======================
class LLMService:
    def __init__(self, model_key: str):
        ident = model_key.split("(")[-1].rstrip(")")
        self.config = AppConfig.MODELS.get(ident, AppConfig.MODELS["local"])

    def generate_cases(self, prompt: str, context: str) -> Dict:
        if self.config.model != "local" and not self.config.api_key:
            return {"status": "error", "message": f"未配置 {self.config.name} API Key"}
        if self.config.model == "local":
            return {"status": "success", "content": self._local_generate(), "message": "本地生成"}

        for attempt in range(AppConfig.API_MAX_RETRIES):
            try:
                client = OpenAI(api_key=self.config.api_key, base_url=self.config.base_url, timeout=AppConfig.API_TIMEOUT)
                messages = [
                    {"role": "system", "content": self._system_prompt()},
                    {"role": "user", "content": f"需求：{prompt}\n\n知识库：\n{context}"}
                ]
                with ThreadPoolExecutor(max_workers=1) as executor:
                    future = executor.submit(
                        lambda: client.chat.completions.create(
                            model=self.config.model, messages=messages, temperature=0.3, max_tokens=2000
                        )
                    )
                    resp = future.result(timeout=AppConfig.API_TIMEOUT)
                    return {"status": "success", "content": resp.choices[0].message.content.strip(), "message": "生成成功"}
            except FuturesTimeoutError:
                if attempt < AppConfig.API_MAX_RETRIES - 1:
                    time.sleep(AppConfig.API_RETRY_DELAY)
                    continue
                return {"status": "error", "message": "请求超时，请稍后重试"}
            except Exception as e:
                if attempt == AppConfig.API_MAX_RETRIES - 1:
                    return {"status": "error", "message": f"生成失败: {str(e)}"}
                time.sleep(AppConfig.API_RETRY_DELAY)
        return {"status": "error", "message": "未知错误"}

    @staticmethod
    def _system_prompt():
        return """你是一名资深测试工程师。请根据给定的需求和知识库，生成规范、可执行的测试用例。输出需严格遵循以下格式，不输出任何额外解释或标记。

    ## 输出格式
    用例ID： TC001
    RDM单号： DEMO-001
    用例名称： [场景]下执行[操作]应[预期]，15字以内，如：用户未勾选协议时点击注册按钮，应提示“请同意用户协议”；
    前置条件： 只写“必须满足才可执行”的状态，不写操作步骤、不重复显而易见的常识、不用完整句子，如已注册账号，浏览器清空缓存，若无则填“无”。不要写入测试数据，言简意赅10字以内
    测试步骤： 1. 具体操作；2. 具体操作；……
    预期结果： 1. 对应步骤1的验证点；2. 对应步骤2的验证点；……

    多个用例之间用空行分隔。

    ## 字段详细规则

    ### 用例名称
    - 简明扼要，直接说明在什么模块或页面或创建下需要验证的测试要点，不超过15个汉字或字符。
    - 格式：简短一句话[模块/场景]下执行[操作]应[预期/校验点],（如：用户未勾选协议时点击注册按钮，应提示“请同意用户协议”；如登录时输入密码错误，提示错误信息并清空密码框）。

    ### 前置条件
    - 只写**必要**的环境、状态或数据前提（如已注册账号，浏览器清空缓存；如订单状态为待支付，创建<1h；）。
    - 如果没有任何必要性前提，必须填“无”。
    - 禁止写入测试步骤中才会使用的具体测试数据（如“用户名=test，密码=123”应写在步骤里）。
    - 去掉所有“已经”、“需要”、“请确保”等虚词
    - 不包含“打开”、“点击”、“输入”等操作动词（那是步骤的内容）

    ### 测试步骤 & 预期结果（一一对应）
    - 步骤与结果的数量必须相同，按相同数字编号一一对应。
    - 若某步骤无预期结果（如中间过渡操作），预期结果对应编号写“无”。
    - **步骤编写要求**：
      - 用“->”表示页面/弹窗/菜单的切换路径（例如：登录页->首页->设置中心）。
      - 产品中的页面名称、按钮文字、提示语、弹窗标题等必须使用双引号（例如：点击“登录”按钮；页面跳转至“个人中心”）。
      - 描述清晰，包含具体操作数据（如输入值、点击坐标/元素）。
    - **预期结果编写要求**：
      1. 优先依据需求说明中描述的操作结果。
      2. 若需求未说明，参考同类成熟产品的典型行为。
      3. 结果必须**肯定无疑义、可客观判定**（如“页面弹出提示‘密码错误’；停留在当前页面”），禁止模糊描述（如“系统正常”）。

    ## 正确示例（请严格模仿此格式）

    用例ID： TC001
    RDM单号： DEMO-001
    用例名称： 登录时输入正确账号密码，跳转首页
    前置条件： 已注册未登录，浏览器无缓存
    测试步骤： 1. 打开“登录”页面；2. 输入用户名“test01”，输入密码“123456”；3. 点击“登录”按钮
    预期结果： 1. 页面显示账号/密码输入框和“登录”按钮；2. 输入框正常接收字符；3. 页面跳转至“首页”，右上角显示“test01”

    用例ID： TC002
    RDM单号： DEMO-002
    用例名称： 登录时用户名密码不一致，提交失败并提示
    前置条件： 已注册用户“test01”，正确密码“123456”
    测试步骤： 1. 进入“登录”页面；2. 输入用户名“test01”，输入密码“654321”；3. 点击“登录”按钮
    预期结果： 1. 页面正常显示；2. 输入框接收输入；3. 页面弹出提示“用户名或密码错误”，不跳转，仍停留在“登录”页面

    用例ID： TC003
    RDM单号： DEMO-003
    用例名称： 回收站批量删除上限为100个文件
    前置条件： “回收站”中至少有101个文件
    测试步骤： 1. 进入“文件管理”页面；2. 全选所有文件（共101个）；3. 点击“批量删除”按钮
    预期结果： 1. 页面显示文件列表；2. 所有文件被勾选；3. 页面提示“单次最多删除100个文件”，无文件被删除

    ## 重要禁止事项
    - 不要输出任何额外的解释、标记或Markdown代码块（如```）。
    - 不要编造需求中未提及的功能或交互细节。
    - 不要在前置条件中写入测试数据（数据必须放在步骤里）。
    - 不要使步骤和结果的数量不一致。
    - 用例名称不要超过15字，不要使用模糊标题（如“测试删除”）。
    - 前置条件不要超过10字，不要使用句子用短语（如“未登录”、“库存>0”）。
    现在，请根据以上规则生成测试用例。"""

    @staticmethod
    def _local_generate():
        return """用例ID： TC001
RDM单号： DEMO-001
用例名称： 功能验证
前置条件： 环境正常
测试步骤： 1.操作;2.验证
预期结果： 1.成功;2.正常"""


# ====================== 用例解析与导出（保持不变） ======================
class TestCaseService:
    @staticmethod
    def parse(content: str, rdm_codes: List[str]) -> List[Dict]:
        cases = []
        blocks = re.split(r'\n\s*\n', content.strip())
        rdm_idx = 0
        for block in blocks:
            if not any(k in block for k in ["用例ID", "用例名称", "测试步骤"]):
                continue
            case = {"RDM单号": "", "用例ID": "", "用例名称": "", "前置条件": "", "测试步骤": "", "预期结果": ""}
            patterns = {
                "用例ID": r"用例ID[:：]\s*([^\n]+)",
                "RDM单号": r"RDM单号[:：]\s*([^\n]+)",
                "用例名称": r"用例名称[:：]\s*([^\n]+)",
                "前置条件": r"前置条件[:：]\s*([^\n]+(?:\n(?!用例|rdm|测试|预期)[^\n]+)*)",
                "测试步骤": r"测试步骤[:：]\s*([^\n]+(?:\n(?!用例|rdm|前置|预期)[^\n]+)*)",
                "预期结果": r"预期结果[:：]\s*([^\n]+(?:\n(?!用例|rdm|前置|测试)[^\n]+)*)"
            }
            for field, pat in patterns.items():
                m = re.search(pat, block, re.DOTALL)
                if m:
                    value = m.group(1).strip()
                    if field in ["用例ID", "RDM单号"]:
                        value = re.sub(r'^\[|\]$', '', value)
                    case[field] = re.sub(r'\n+', ' ', value)
            if not case["RDM单号"] and rdm_codes:
                case["RDM单号"] = rdm_codes[rdm_idx % len(rdm_codes)]
                rdm_idx += 1
            if case["用例ID"]:
                cases.append(case)
        return cases

class ExportService:
    @staticmethod
    def to_csv(cases: List[Dict]) -> bytes:
        df = pd.DataFrame(cases)
        cols = ["RDM单号", "用例ID", "用例名称", "前置条件", "测试步骤", "预期结果"]
        df = df[[c for c in cols if c in df.columns]]
        buf = io.BytesIO()
        df.to_csv(buf, index=False, encoding='utf-8-sig')
        return buf.getvalue()


# ====================== 主视图（完全不需要改动） ======================
class MainView:
    def __init__(self):
        if "current_project" not in st.session_state:
            projects = ProjectManager.get_all_projects()
            if projects:
                st.session_state.current_project = projects[0]
            else:
                ProjectManager.create_project("默认项目")
                st.session_state.current_project = "默认项目"
        if "uploader_key" not in st.session_state:
            st.session_state.uploader_key = 0
        if "show_new_project_input" not in st.session_state:
            st.session_state.show_new_project_input = False

        self.current_project = st.session_state.current_project
        self.kb = EnhancedKnowledgeBase(self.current_project)

    def render(self):
        st.title("测试用例智构系统")

        with st.sidebar:
            st.subheader("模型选择")
            model_key = st.selectbox("", AppConfig.get_model_list(), label_visibility="collapsed")
            st.divider()

            col_title, col_new = st.columns([3, 3])
            with col_title:
                st.subheader("项目管理")
            with col_new:
                if st.button("+ 新建项目", key="new_project_btn", use_container_width=True):
                    st.session_state.show_new_project_input = True

            if st.session_state.show_new_project_input:
                new_project_name = st.text_input("项目名称", placeholder="请输入项目名称", key="dialog_project_name")
                col_ok, col_cancel = st.columns(2)
                with col_ok:
                    if st.button("确定", key="confirm_new_project"):
                        if new_project_name and new_project_name.strip():
                            if ProjectManager.create_project(new_project_name):
                                st.success(f"项目 '{new_project_name}' 创建成功")
                                st.session_state.current_project = new_project_name
                                st.session_state.show_new_project_input = False
                                st.rerun()
                            else:
                                st.error("项目创建失败，可能名称已存在")
                        else:
                            st.warning("请输入项目名称")
                with col_cancel:
                    if st.button("取消", key="cancel_new_project", use_container_width=True):
                        st.session_state.show_new_project_input = False
                        st.rerun()

            projects = ProjectManager.get_all_projects()
            if projects:
                selected_project = st.selectbox(
                    "当前项目",
                    options=projects,
                    index=projects.index(self.current_project) if self.current_project in projects else 0,
                    key="project_selector"
                )
                if selected_project != self.current_project:
                    st.session_state.current_project = selected_project
                    st.rerun()
            else:
                st.warning("暂无项目，请先创建")

            st.divider()

            st.subheader("文档上传")
            uploaded_files = st.file_uploader(
                "选择文件（可多选）",
                type=list(AppConfig.ALLOWED_FILE_TYPES),
                accept_multiple_files=True,
                key=f"file_uploader_{st.session_state.uploader_key}",
                help="支持Word/PDF/Excel/PPT/Xmind等格式，图片OCR识别开发中",
            )
            if st.button("上传", key="upload_btn", use_container_width=True):
                if uploaded_files:
                    success_count = 0
                    fail_count = 0
                    for uploaded_file in uploaded_files:
                        success = self.kb.upload_file(uploaded_file.name, uploaded_file.getvalue())
                        if success:
                            success_count += 1
                        else:
                            fail_count += 1
                    if success_count > 0:
                        st.success(f"成功上传 {success_count} 个文件")
                    if fail_count > 0:
                        st.warning(f"{fail_count} 个文件已存在，未重复上传")
                    st.session_state.uploader_key += 1
                    time.sleep(1)
                    st.rerun()
                else:
                    st.warning("请先选择文件")

        self._kb_panel()
        st.divider()
        self._gen_panel(model_key)

    def _kb_panel(self):
        st.subheader(f"知识库 - {self.current_project}")
        files = self.kb.get_file_list()
        built_files = self.kb.get_built_files_safe()
        if files:
            for f in files:
                col1, col2, col3, _ = st.columns([0.5, 3, 0.5, 2])
                is_built = f['name'] in built_files
                if is_built:
                    col1.write("✅已构建")
                else:
                    col1.write("⏳未构建")
                col2.write(f"📄 {f['name']} ({f['size']})")
                delete_key = f"del_{self.current_project}_{f['name']}"
                if is_built:
                    with col3.popover("删除", use_container_width=False):
                        st.write(f"确定要删除文档 **{f['name']}** 吗？")
                        st.write("删除后，该文档构建的索引也会被清除。")
                        col_confirm, col_cancel = st.columns(2)
                        with col_confirm:
                            if st.button("确认删除", key=f"confirm_{delete_key}", type="primary"):
                                self._delete_file_and_rebuild(f['name'])
                                st.rerun()
                else:
                    if col3.button("删除", key=delete_key):
                        self._delete_file_and_rebuild(f['name'])
                        st.rerun()
        else:
            st.info("暂无文档，请先上传文件")

        col_btn1, col_btn2, _ = st.columns([1, 1, 6])
        with col_btn1:
            if st.button("构建知识库", type="primary", use_container_width=True):
                if not files:
                    st.warning("请先上传文件")
                else:
                    with st.status("正在构建知识库...", expanded=True) as status:
                        prog = st.progress(0)
                        stat = st.empty()
                        def cb(p, msg):
                            prog.progress(p)
                            stat.text(msg)
                        res = self.kb.build_knowledge_base(cb, force=False)
                        if res["status"] == "success":
                            status.update(state="complete")
                            st.toast(res['message'], icon="✅")
                        else:
                            status.update(label="", state="error")
                            st.toast(res['message'], icon="❌")
                    time.sleep(1)
                    st.rerun()
        with col_btn2:
            if st.button("刷新索引", use_container_width=True):
                res = self.kb.refresh_index()
                if res["status"] == "success":
                    st.toast(res["message"], icon="🔄")
                else:
                    st.toast(res["message"], icon="⚠️")
                time.sleep(1)
                st.rerun()

    def _delete_file_and_rebuild(self, filename: str):
        self.kb.delete_file(filename)
        st.rerun()

    def _gen_panel(self, model_key):
        st.subheader("生成测试用例")
        input_type = st.radio("输入方式", ["文本输入", "RDM单号"], horizontal=True)
        if input_type == "文本输入":
            prompt = st.text_area("需求描述", height=150, key="prompt_text")
        else:
            rdm = st.text_input("RDM单号", key="rdm_input")
            prompt = st.text_area("需求描述", value=rdm, height=150, key="prompt_with_rdm") if rdm else ""

        col_check, _ = st.columns([1, 6])
        with col_check:
            use_rag = st.checkbox("启用RAG知识库", value=True,
                                  help="从当前项目的知识库中检索相关内容，提升生成质量")

        col_gen, col_clear, _ = st.columns([1, 1, 6])
        with col_gen:
            generate_clicked = st.button("生成测试用例", type="primary", use_container_width=True)
        with col_clear:
            clear_clicked = st.button("清空结果", use_container_width=True)

        if clear_clicked:
            if "cases" in st.session_state:
                del st.session_state.cases
            st.rerun()

        if generate_clicked:
            if not prompt:
                st.warning("请输入需求描述")
                return

            with st.spinner("生成中..."):
                context = ""
                if use_rag:
                    files = self.kb.get_file_list()
                    if files:
                        if not self.kb.index_loaded:
                            self.kb.refresh_index()
                        if self.kb.index_loaded:
                            with st.expander("RAG检索详情"):
                                results = self.kb.search_knowledge(prompt, top_k=5)
                                if results:
                                    for r in results:
                                        st.write(f"相似度 {r['similarity']:.3f} | 来源 {r['metadata'].get('source','')}")
                                        st.caption(r['content'][:200])
                                else:
                                    st.info("未检索到相关知识")
                            context = self.kb.get_knowledge_context(prompt, max_chunks=5)
                        else:
                            st.info("知识库尚未构建，将直接使用模型生成。如需检索知识，请先构建知识库。")
                    else:
                        st.info("当前项目无文档，将直接使用模型生成。")

                llm = LLMService(model_key)
                resp = llm.generate_cases(prompt, context)

                if resp["status"] == "error":
                    st.error(resp["message"])
                else:
                    st.success(resp["message"])
                    rdm_codes = re.findall(r'[A-Za-z0-9_]+-\d+', prompt)
                    cases = TestCaseService.parse(resp["content"], rdm_codes)
                    if cases:
                        st.session_state.cases = cases
                        st.success(f"解析出 {len(cases)} 个测试用例")
                    else:
                        st.warning("解析失败，显示原始内容")
                        st.code(resp["content"])

        if "cases" in st.session_state:
            self._show_results()

    @staticmethod
    def _show_results():
        cases = st.session_state.cases
        if not cases:
            return
        st.subheader("测试用例")
        df = pd.DataFrame(cases)
        cols = ["RDM单号", "用例ID", "用例名称", "前置条件", "测试步骤", "预期结果"]
        df = df[[c for c in cols if c in df.columns]]
        st.markdown("""
        <style>
        .auto-wrap-table { width: 100%; border-collapse: collapse; }
        .auto-wrap-table th, .auto-wrap-table td { border: 1px solid #ddd; padding: 8px; text-align: left; vertical-align: top; white-space: normal; word-wrap: break-word; }
        .auto-wrap-table th:nth-child(1), .auto-wrap-table td:nth-child(1) { width: 100px; }
        .auto-wrap-table th:nth-child(2), .auto-wrap-table td:nth-child(2) { width: 80px; }
        .auto-wrap-table th:nth-child(3), .auto-wrap-table td:nth-child(3) { width: 220px; }
        .auto-wrap-table th:nth-child(4), .auto-wrap-table td:nth-child(4) { width: 160px; }
        .auto-wrap-table th:nth-child(5), .auto-wrap-table td:nth-child(5) { width: 320px; }
        .auto-wrap-table th:nth-child(6), .auto-wrap-table td:nth-child(6) { width: 320px; }
        </style>
        """, unsafe_allow_html=True)
        html_table = df.to_html(index=False, classes='auto-wrap-table', escape=False)
        st.markdown(html_table, unsafe_allow_html=True)
        csv = ExportService.to_csv(cases)
        st.download_button(
            "导出 CSV",
            data=csv,
            file_name=f"cases_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv",
            mime="text/csv",
            use_container_width=True
        )


# ====================== 入口 ======================
def main():
    # 本地开发时确保临时目录存在（实际上不需要创建 OSS 相关目录）
    MainView().render()

if __name__ == "__main__":
    main()