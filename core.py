import streamlit as st
import os
import re
import json
import pickle
import io
import warnings
import tempfile
from pathlib import Path
from datetime import datetime
from typing import List, Dict, Optional, Callable
import numpy as np
import pandas as pd
import oss2
from oss2.exceptions import NoSuchKey
import requests
from bs4 import BeautifulSoup
import docx
import PyPDF2
from openpyxl import load_workbook
from pptx import Presentation
import jieba
from config import AppConfig
warnings.filterwarnings('ignore')

# ====================== 阿里云 OSS 辅助函数(获取配置） ======================
@st.cache_resource
def get_oss_bucket():
    """从 st.secrets 获取 OSS 配置并返回 Bucket 对象（缓存）"""
    try:
        auth = oss2.Auth(
            st.secrets["OSS_ACCESS_KEY_ID"],
            st.secrets["OSS_ACCESS_KEY_SECRET"]
        )
        bucket = oss2.Bucket(
            auth,
            st.secrets["OSS_ENDPOINT"],
            st.secrets["OSS_BUCKET_NAME"]
        )
        return bucket
    except Exception as e:
        st.error(f"OSS 配置错误: {e}")
        st.stop()

# ====================== 项目管理（基于 OSS 前缀） ======================
class ProjectManager:
    @staticmethod
    @st.cache_data(ttl=60)  # 缓存60秒，减少OSS调用
    def get_all_projects() -> List[str]:
        """通过列出 OSS 前缀获取所有项目名称"""
        bucket = get_oss_bucket()
        prefix = f"{AppConfig.BASE_ROOT}/"
        try:
            projects = []
            for obj in oss2.ObjectIterator(bucket, prefix=prefix, delimiter='/'):
                # obj.key 格式如 "knowledge_projects/项目名/"
                if obj.key.endswith('/'):
                    proj_name = obj.key[len(prefix):-1]
                    if proj_name:
                        projects.append(proj_name)
            return sorted(projects)
        except Exception:
            return []

    @staticmethod
    def create_project(project_name: str) -> bool:
        if not project_name or not project_name.strip():
            return False
        project_name = project_name.strip()
        if not re.match(r'^[\u4e00-\u9fa5a-zA-Z0-9_-]+$', project_name):
            st.error("项目名称只能包含中文、字母、数字、下划线和中划线")
            return False
        existing = ProjectManager.get_all_projects()
        if project_name in existing:
            return False
        bucket = get_oss_bucket()
        marker_key = f"{AppConfig.BASE_ROOT}/{project_name}/.project"
        try:
            bucket.put_object(marker_key, b'')
            bucket.put_object(f"{AppConfig.BASE_ROOT}/{project_name}/knowledge_base/.keep", b'')
            bucket.put_object(f"{AppConfig.BASE_ROOT}/{project_name}/vector_store/.keep", b'')
            # 清除缓存
            st.cache_data.clear()
            return True
        except Exception:
            return False

    @staticmethod
    def get_project_path(project_name: str) -> str:
        return f"{AppConfig.BASE_ROOT}/{project_name}"

    @staticmethod
    def get_kb_path(project_name: str) -> str:
        return f"{AppConfig.BASE_ROOT}/{project_name}/knowledge_base"

    @staticmethod
    def get_vector_store_path(project_name: str) -> str:
        return f"{AppConfig.BASE_ROOT}/{project_name}/vector_store"


# ====================== 知识库封装（OSS 存储） ======================
class EnhancedKnowledgeBase:
    def __init__(self, project_name: str):
        self.project_name = project_name
        self.oss_bucket = get_oss_bucket()
        self.kb_prefix = ProjectManager.get_kb_path(project_name)  # OSS 前缀
        self.vector_prefix = ProjectManager.get_vector_store_path(project_name)

        # 临时目录 - 使用项目名区分
        self.temp_dir = Path(tempfile.gettempdir()) / f"kb_{project_name}"
        self.temp_dir.mkdir(parents=True, exist_ok=True)

        self.vector_store = SimpleVectorStore(
            str(self.temp_dir / "vector_store"),
            self.oss_bucket,
            self.vector_prefix
        )
        self.chunk_size = 800
        self.chunk_overlap = 200
        self.index_loaded = False
        self._file_list_cache = None  # 文件列表缓存
        self._auto_load()

    def _auto_load(self):
        """自动加载索引，不进行额外的有效性检查以避免重复OSS调用"""
        try:
            files = self.get_file_list(with_metadata=True)
            if files:
                # 直接尝试加载索引，让load_index自己判断
                if self.vector_store.load_index():
                    self.index_loaded = True
        except Exception as e:
            # 静默失败，后续操作会处理
            pass

    def get_file_list(self, with_metadata=False):
        """获取文件列表 - 不使用缓存，直接读取OSS"""
        files = []
        prefix = f"{self.kb_prefix}/"
        try:
            for obj in oss2.ObjectIterator(self.oss_bucket, prefix=prefix):
                key = obj.key
                filename = key[len(prefix):]
                if not filename or filename.startswith('.') or filename.endswith('/'):
                    continue
                ext = filename.split('.')[-1].lower()
                if ext in AppConfig.ALLOWED_FILE_TYPES:
                    size_kb = obj.size / 1024
                    info = {"name": filename, "size": f"{size_kb:.1f}KB"}
                    if with_metadata:
                        if hasattr(obj.last_modified, 'timestamp'):
                            mtime = obj.last_modified.timestamp()
                        else:
                            mtime = float(obj.last_modified)
                        info.update({"size_bytes": obj.size, "mtime": mtime})
                    files.append(info)
            return sorted(files, key=lambda x: x['name'])
        except Exception as e:
            st.error(f"OSS 列表失败: {e}")
            return []

    def upload_file(self, filename: str, content: bytes) -> bool:
        key = f"{self.kb_prefix}/{filename}"
        try:
            # 检查是否存在
            self.oss_bucket.head_object(key)
            return False
        except NoSuchKey:
            pass
        try:
            self.oss_bucket.put_object(key, content)
            self.index_loaded = False
            return True
        except Exception:
            return False

    def delete_file(self, filename: str):
        key = f"{self.kb_prefix}/{filename}"
        try:
            self.oss_bucket.delete_object(key)
        except Exception:
            pass
        self.vector_store.remove_file(filename)
        remaining_files = self.get_file_list(with_metadata=True)
        self.vector_store.update_metadata(remaining_files)
        self.index_loaded = len(remaining_files) > 0 and self.vector_store.load_index()

    #========================= 文档解析（从 OSS 读取字节流） ====================#
    def _read_file_from_oss(self, filename: str) -> bytes:
        key = f"{self.kb_prefix}/{filename}"
        return self.oss_bucket.get_object(key).read()

    def _extract_text_from_bytes(self, filename: str, content_bytes: bytes) -> str:
        ext = filename.split('.')[-1].lower()
        # 对于需要文件路径的解析器，写入临时文件
        if ext in ['docx', 'doc', 'pdf', 'xlsx', 'xls', 'pptx', 'xmind']:
            with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as tmp:
                tmp.write(content_bytes)
                tmp_path = tmp.name
            try:
                return self._extract_text_from_path(tmp_path, ext)
            finally:
                os.unlink(tmp_path)
        else:
            # txt, csv, json, md 直接解码
            if ext in ['txt', 'md', 'json']:
                return content_bytes.decode('utf-8', errors='ignore')
            elif ext == 'csv':
                df = pd.read_csv(io.BytesIO(content_bytes))
                return df.to_string()
            else:
                return f"暂无法识别该类型 (.{ext})，请联系管理员"

    def _extract_text_from_path(self, filepath: str, ext: str) -> str:
        parsers = {
            "docx": lambda p: '\n'.join([para.text for para in docx.Document(p).paragraphs]),
            "doc": lambda p: '\n'.join([para.text for para in docx.Document(p).paragraphs]),
            "pdf": lambda p: '\n'.join([page.extract_text() for page in PyPDF2.PdfReader(p).pages]),
            "xlsx": lambda p: self._parse_excel(p),
            "xls": lambda p: self._parse_excel(p),
            "pptx": lambda p: '\n'.join(
                [shape.text for slide in Presentation(p).slides for shape in slide.shapes if hasattr(shape, "text")]),
            "xmind": lambda p: self._parse_xmind(p),
        }
        if ext in parsers:
            try:
                return parsers[ext](filepath)
            except Exception as e:
                return f"提取失败: {e}"
        return f"暂无法识别该类型 (.{ext})，请联系管理员"

    @staticmethod
    def _parse_excel(filepath: str) -> str:  # Excel 表格完整转成纯文本
        wb = load_workbook(filepath, data_only=True)
        all_text = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = ' '.join([str(cell) for cell in row if cell is not None])
                if row_text:
                    all_text.append(row_text)
        return '\n'.join(all_text)

    @staticmethod
    @staticmethod
    def _parse_xmind(filepath: str) -> str:
        """解析XMind文件（兼容所有版本，直接解压读取JSON）"""
        all_texts = []
        try:
            import zipfile
            import json

            with zipfile.ZipFile(filepath, 'r') as zf:
                # 优先处理 content.json
                if 'content.json' in zf.namelist():
                    data = json.load(zf.open('content.json'))
                    # 处理 data 可能是列表或字典的情况
                    sheets_data = []
                    if isinstance(data, dict):
                        sheets_data = data.get('sheets', [])
                    elif isinstance(data, list):
                        # 有些版本直接返回 sheets 列表
                        sheets_data = data
                    else:
                        return "无法识别的XMind JSON结构"

                    def extract_from_topic(topic):
                        texts = []
                        if not isinstance(topic, dict):
                            return texts
                        title = topic.get('title')
                        if title and isinstance(title, str) and title.strip():
                            texts.append(title.strip())
                        # 备注
                        notes = topic.get('notes')
                        if isinstance(notes, dict):
                            plain = notes.get('plain')
                            if isinstance(plain, dict):
                                content = plain.get('content')
                                if content and isinstance(content, str) and content.strip():
                                    texts.append(content.strip())
                        # 子主题
                        children = topic.get('children')
                        if isinstance(children, dict):
                            attached = children.get('attached', [])
                            if isinstance(attached, list):
                                for child in attached:
                                    texts.extend(extract_from_topic(child))
                        return texts

                    for sheet in sheets_data:
                        if not isinstance(sheet, dict):
                            continue
                        sheet_title = sheet.get('title')
                        if sheet_title and isinstance(sheet_title, str) and sheet_title.strip():
                            all_texts.append(sheet_title.strip())
                        root_topic = sheet.get('rootTopic')
                        if isinstance(root_topic, dict):
                            all_texts.extend(extract_from_topic(root_topic))

                # 兼容旧版 content.xml
                elif 'content.xml' in zf.namelist():
                    content = zf.read('content.xml').decode('utf-8')
                    # 提取 <title> 和 <plain-text> 标签内容
                    title_matches = re.findall(r'<title[^>]*>([^<]+)</title>', content)
                    plain_matches = re.findall(r'<plain-text[^>]*>([^<]+)</plain-text>', content)
                    all_texts.extend(title_matches)
                    all_texts.extend(plain_matches)
                else:
                    return "无法识别 XMind 文件格式，缺少 content.json 或 content.xml"

        except Exception as e:
            return f"XMind解析失败: {str(e)}"

        # 去重并返回
        seen = set()
        unique_texts = []
        for text in all_texts:
            if text and text not in seen:
                seen.add(text)
                unique_texts.append(text)
        return '\n'.join(unique_texts) if unique_texts else "未提取到有效文本内容"

    #=============================== 分块与构建 =========================#
    def _chunk_text(self, text: str, source: str) -> List[Dict]:
        prefixed_text = f"[文件：{source}]\n{text}"
        chunks = []
        paragraphs = prefixed_text.split('\n')
        current = ""
        chunk_id = 0
        for para in paragraphs:
            para = para.strip()
            if not para:
                continue
            if len(para) > self.chunk_size:
                sentences = re.split(r'[。！？；]', para)
                for sentence in sentences:
                    sentence = sentence.strip()
                    if not sentence:
                        continue
                    if len(current) + len(sentence) <= self.chunk_size:
                        current += sentence + "。"
                    else:
                        if current:
                            chunks.append(
                                {"content": current.strip(), "metadata": {"source": source, "chunk_id": chunk_id}})
                            chunk_id += 1
                        overlap = current[-self.chunk_overlap:] if current else ""
                        current = overlap + sentence + "。"
            else:
                if len(current) + len(para) <= self.chunk_size:
                    current += para + "\n"
                else:
                    if current:
                        chunks.append(
                            {"content": current.strip(), "metadata": {"source": source, "chunk_id": chunk_id}})
                        chunk_id += 1
                    overlap = current[-self.chunk_overlap:] if current else ""
                    current = overlap + para + "\n"
        if current:
            chunks.append({"content": current.strip(), "metadata": {"source": source, "chunk_id": chunk_id}})
        return chunks

    def build_knowledge_base(self, progress_callback: Optional[Callable] = None, force=False) -> Dict:
        try:
            files = self.get_file_list(with_metadata=True)
            if not files:
                return {"status": "error", "message": "无文件"}

            if self.index_loaded:
                built_files = self.vector_store.get_built_files()
            else:
                if self.vector_store.load_index():
                    self.index_loaded = True
                    built_files = self.vector_store.get_built_files()
                else:
                    built_files = set()

            files_to_build = [f for f in files if f["name"] not in built_files]

            if not force and not files_to_build:
                stats = self.vector_store.get_stats()
                return {"status": "success", "message": f"所有文档已构建（共 {stats['documents']} 个知识片段）",
                        "chunks": stats['documents']}

            # 增量构建（跳过已构建文档）
            if built_files and not force:
                new_chunks = []
                total_to_build = len(files_to_build)
                for i, f in enumerate(files_to_build):
                    if progress_callback:
                        progress_callback(i / total_to_build, f"处理 {f['name']} (解析中...)")
                    content_bytes = self._read_file_from_oss(f["name"])
                    text = self._extract_text_from_bytes(f["name"], content_bytes)
                    if text.startswith("请安装") or text.startswith("暂无法识别") or text.startswith("提取失败"):
                        continue
                    chunks = self._chunk_text(text, f["name"])
                    new_chunks.extend(chunks)
                    if progress_callback:
                        progress_callback(i / total_to_build, f"处理 {f['name']} → 生成 {len(chunks)} 个片段")

                if not new_chunks:
                    return {"status": "error", "message": "没有可构建的有效内容，请检查文件格式"}

                all_docs = self.vector_store.documents + [
                    {"id": i + len(self.vector_store.documents), "text": c["content"], "metadata": c["metadata"]}
                    for i, c in enumerate(new_chunks)
                ]
                texts = [doc["text"] for doc in all_docs]
                self.vector_store._build_vocab_and_idf(texts)
                vectors = [self.vector_store._text_to_vector(t) for t in texts]
                self.vector_store.vectors = np.array(vectors)
                self.vector_store.documents = all_docs
                # 保存到本地并上传 OSS
                with open(self.vector_store.index_file, 'w', encoding='utf-8') as f:
                    json.dump({
                        "documents": self.vector_store.documents,
                        "vocab": self.vector_store.vocab,
                        "idf": self.vector_store.idf
                    }, f, ensure_ascii=False)
                np.save(self.vector_store.vectors_file, self.vector_store.vectors)
                self.vector_store._upload_to_oss(self.vector_store.index_file, self.vector_store._oss_key("index.json"))
                self.vector_store._upload_to_oss(self.vector_store.vectors_file,
                                                 self.vector_store._oss_key("vectors.npy"))
                files_info = [{"name": f["name"], "size_bytes": f["size_bytes"], "mtime": f["mtime"]} for f in files]
                self.vector_store.update_metadata(files_info)
                self.index_loaded = True
                return {"status": "success",
                        "message": f"增量构建成功，新增 {len(files_to_build)} 个文件，{len(new_chunks)} 个知识片段",
                        "files": len(files), "chunks": len(all_docs)}
            else:
                # 全量构建
                all_chunks = []
                missing_deps = set()
                unsupported_files = []
                total_files = len(files)
                for i, f in enumerate(files):
                    if progress_callback:
                        progress_callback(i / total_files, f"处理 {f['name']} (解析中...)")
                    content_bytes = self._read_file_from_oss(f["name"])
                    text = self._extract_text_from_bytes(f["name"], content_bytes)
                    if text.startswith("请安装"):
                        missing_deps.add(text)
                        continue
                    if text.startswith("暂无法识别该类型"):
                        unsupported_files.append(f"{f['name']}: {text}")
                        continue
                    if text.startswith("提取失败"):
                        unsupported_files.append(f"{f['name']}: {text}")
                        continue
                    chunks = self._chunk_text(text, f["name"])
                    all_chunks.extend(chunks)
                    if progress_callback:
                        progress_callback(i / total_files, f"处理 {f['name']} → 生成 {len(chunks)} 个片段")
                if missing_deps:
                    st.warning("\n".join(missing_deps))
                if unsupported_files:
                    for err in unsupported_files:
                        st.error(err)
                    if not all_chunks:
                        return {"status": "error", "message": "所有文件均无法解析，请检查依赖或文件格式"}
                if progress_callback:
                    progress_callback(0.95, "构建向量索引...")
                files_info = [{"name": f["name"], "size_bytes": f["size_bytes"], "mtime": f["mtime"]} for f in files]
                ok = self.vector_store.build_index(all_chunks, files_info)
                if ok:
                    self.index_loaded = True
                    return {"status": "success",
                            "message": f"构建成功，共 {len(files)} 个文件，{len(all_chunks)} 个知识片段",
                            "files": len(files), "chunks": len(all_chunks)}
                else:
                    return {"status": "error", "message": "索引失败"}
        except Exception as e:
            st.error(f"构建知识库时发生未预期的错误: {str(e)}")
            return {"status": "error", "message": str(e)}

    def search_knowledge(self, query: str, top_k=5):
        if not self.index_loaded and not self.vector_store.load_index():
            return []
        return self.vector_store.search(query, top_k)

    def get_knowledge_context(self, query: str, max_chunks=5) -> str:
        results = self.search_knowledge(query, max_chunks)
        if not results:
            return ""
        parts = []
        for i, r in enumerate(results, 1):
            parts.append(
                f"【知识片段 {i}】来源：{r['metadata'].get('source', '未知')} 相关度：{r['similarity']:.2%}\n{r['content']}")
        return "\n\n---\n\n".join(parts)

    def get_built_files_safe(self) -> set:
        if self.index_loaded:
            return self.vector_store.get_built_files()
        else:
            if self.vector_store.load_index():
                self.index_loaded = True
                return self.vector_store.get_built_files()
            return set()

    def refresh_index(self) -> Dict:
        """刷新索引，清除缓存"""
        self.index_loaded = False
        files = self.get_file_list(with_metadata=True)
        if not files:
            self.vector_store.clear_index()
            return {"status": "error", "message": "知识库无文件"}
        if self.vector_store.load_index():
            self.index_loaded = True
            stats = self.vector_store.get_stats()
            built_files = self.vector_store.get_built_files()
            all_files = {f["name"] for f in files}
            pending = all_files - built_files
            if pending:
                msg = f"已加载索引，包含 {stats['documents']} 个片段。注意：有 {len(pending)} 个文件未构建，请点击「构建知识库」增量添加。"
            else:
                msg = f"已加载索引，包含 {stats['documents']} 个知识片段"
            return {"status": "success", "message": msg}
        return {"status": "error", "message": "未找到有效索引，请先构建知识库"}


# ====================== 轻量级 TF‑IDF 向量检索（RAG检索，支持OSS） ======================
class SimpleVectorStore:
    def __init__(self, store_dir: str, oss_bucket, oss_prefix: str):
        """
        store_dir: 本地临时目录，用于存放下载的索引文件
        oss_bucket: OSS Bucket 对象
        oss_prefix: OSS 中该项目的向量存储前缀
        """
        self.store_dir = store_dir
        self.oss_bucket = oss_bucket
        self.oss_prefix = oss_prefix
        Path(store_dir).mkdir(parents=True, exist_ok=True)

        self.index_file = os.path.join(store_dir, "index.json")
        self.vectors_file = os.path.join(store_dir, "vectors.npy")
        self.metadata_file = os.path.join(store_dir, "metadata.pkl")
        self.documents = []
        self.vocab = {}
        self.idf = {}
        self.vectors = None
        self._cache_valid = False  # 缓存有效性标记

    def _oss_key(self, filename: str) -> str:
        """生成 OSS 对象键"""
        return f"{self.oss_prefix}/{filename}"

    def _download_from_oss(self, local_path: str, object_key: str) -> bool:
        """从 OSS 下载文件到本地，如果不存在返回 False"""
        try:
            self.oss_bucket.get_object_to_file(object_key, local_path)
            return True
        except NoSuchKey:
            return False
        except Exception as e:
            st.warning(f"下载文件失败 {object_key}: {e}")
            return False

    def _upload_to_oss(self, local_path: str, object_key: str):
        """上传本地文件到 OSS"""
        try:
            self.oss_bucket.put_object_from_file(object_key, local_path)
        except Exception as e:
            st.error(f"上传文件失败 {object_key}: {e}")

    def _delete_oss_object(self, object_key: str):
        """删除 OSS 对象"""
        try:
            self.oss_bucket.delete_object(object_key)
        except Exception:
            pass

    @staticmethod
    def _tokenize(text: str) -> List[str]:
        text = text.lower()
        words = jieba.lcut(text)
        return [w for w in words if len(w.strip()) > 1 or w.isalnum()]

    @staticmethod
    def _clean_query(query: str) -> str:
        query = re.sub(r'\b[A-Za-z0-9]+[_-][A-Za-z0-9]+\b', '', query)
        query = re.sub(r'\b[A-Z]{2,}[0-9]+\b', '', query)
        query = re.sub(r'\b[A-Z]+-\d+\b', '', query)
        query = re.sub(r'\b\d{4}[-/]\d{1,2}[-/]\d{1,2}\b', '', query)
        query = re.sub(r'\b\d{1,2}[-/]\d{1,2}[-/]\d{4}\b', '', query)
        query = re.sub(r'\b\d{8}\b', '', query)
        query = re.sub(r'\s+', ' ', query).strip()
        return query

    def _build_vocab_and_idf(self, all_texts: List[str]):
        doc_count = {}
        for text in all_texts:
            unique_words = set(self._tokenize(text))
            for w in unique_words:
                doc_count[w] = doc_count.get(w, 0) + 1

        total_docs = len(all_texts)
        self.vocab = {}
        self.idf = {}
        for w, df in doc_count.items():
            if 0.01 < df / total_docs < 0.8:
                idx = len(self.vocab)
                self.vocab[w] = idx
                self.idf[w] = np.log((total_docs + 1) / (df + 0.5))

    def _text_to_vector(self, text: str) -> np.ndarray:
        words = self._tokenize(text)
        if not words or not self.vocab:
            return np.zeros(len(self.vocab))
        tf = {}
        for w in words:
            if w in self.vocab:
                tf[w] = tf.get(w, 0) + 1
        vec = np.zeros(len(self.vocab))
        for w, cnt in tf.items():
            idx = self.vocab[w]
            tf_val = np.log1p(cnt / len(words))
            vec[idx] = tf_val * self.idf.get(w, 1.0)
        norm = np.linalg.norm(vec)
        if norm > 0:
            vec /= norm
        return vec

    def build_index(self, documents: List[Dict], files_info: List[Dict] = None) -> bool:
        try:
            if not documents:
                return False
            self.documents = [{"id": i, "text": d["content"], "metadata": d.get("metadata", {})}
                              for i, d in enumerate(documents)]
            texts = [d["text"] for d in self.documents]
            self._build_vocab_and_idf(texts)
            vectors = [self._text_to_vector(t) for t in texts]
            self.vectors = np.array(vectors)

            # 保存到本地临时文件
            with open(self.index_file, 'w', encoding='utf-8') as f:
                json.dump({
                    "documents": self.documents,
                    "vocab": self.vocab,
                    "idf": self.idf
                }, f, ensure_ascii=False)
            np.save(self.vectors_file, self.vectors)
            if files_info:
                with open(self.metadata_file, 'wb') as f:
                    pickle.dump({"files_info": files_info, "updated": datetime.now().isoformat()}, f)

            # 上传到 OSS
            self._upload_to_oss(self.index_file, self._oss_key("index.json"))
            self._upload_to_oss(self.vectors_file, self._oss_key("vectors.npy"))
            if files_info and os.path.exists(self.metadata_file):
                self._upload_to_oss(self.metadata_file, self._oss_key("metadata.pkl"))

            self._cache_valid = True
            return True
        except Exception as e:
            st.error(f"构建索引失败: {e}")
            return False

    def load_index(self) -> bool:
        try:
            # 检查本地文件是否有效且最新
            if self._cache_valid and self.vectors is not None:
                return True

            if not self._download_from_oss(self.index_file, self._oss_key("index.json")):
                return False
            if not self._download_from_oss(self.vectors_file, self._oss_key("vectors.npy")):
                return False
            self._download_from_oss(self.metadata_file, self._oss_key("metadata.pkl"))

            with open(self.index_file, 'r', encoding='utf-8') as f:
                data = json.load(f)
                self.documents = data["documents"]
                self.vocab = data["vocab"]
                self.idf = data["idf"]
            self.vectors = np.load(self.vectors_file)
            self._cache_valid = True
            return True
        except Exception as e:
            st.warning(f"加载索引失败: {e}")
            return False

    def search(self, query: str, top_k: int = 5) -> List[Dict]:
        if self.vectors is None or len(self.documents) == 0:
            return []
        cleaned_query = self._clean_query(query)
        if not cleaned_query:
            cleaned_query = query
        q_vec = self._text_to_vector(cleaned_query)
        scores = np.dot(self.vectors, q_vec)
        top_idx = np.argsort(scores)[-top_k:][::-1]
        results = []
        for idx in top_idx:
            sim = float(scores[idx])
            if sim >= AppConfig.RAG_SIMILARITY_THRESHOLD:
                doc = self.documents[idx]
                results.append({
                    "content": doc["text"],
                    "metadata": doc["metadata"],
                    "similarity": sim
                })
        return results

    def is_index_valid(self, current_files_info: List[Dict]) -> bool:
        """检查 OSS 上的元数据文件是否与当前文件列表一致"""
        temp_meta = os.path.join(self.store_dir, "metadata_check.pkl")
        try:
            if not self._download_from_oss(temp_meta, self._oss_key("metadata.pkl")):
                return False
            with open(temp_meta, 'rb') as f:
                meta = pickle.load(f)
            saved = {f["name"]: f for f in meta["files_info"]}
            curr = {f["name"]: f for f in current_files_info}
            if set(saved.keys()) != set(curr.keys()):
                return False
            for name, cf in curr.items():
                sf = saved[name]
                if sf.get("size_bytes") != cf.get("size_bytes") or sf.get("mtime") != cf.get("mtime"):
                    return False
            return True
        except Exception as e:
            # 文件不存在或解析失败都返回 False
            return False
        finally:
            # 清理临时文件
            if os.path.exists(temp_meta):
                try:
                    os.remove(temp_meta)
                except:
                    pass

    def get_built_files(self) -> set:
        if not self.documents:
            return set()
        return {doc["metadata"].get("source") for doc in self.documents if "source" in doc["metadata"]}

    def get_stats(self) -> Dict:
        return {
            "documents": len(self.documents),
            "vocab_size": len(self.vocab),
            "has_index": self.vectors is not None
        }

    def remove_file(self, filename: str) -> bool:
        """从当前内存索引中删除指定文件的分块，并重建索引，然后上传到 OSS"""
        if not self.documents:
            return False
        new_docs = [doc for doc in self.documents if doc["metadata"].get("source") != filename]
        if len(new_docs) == len(self.documents):
            return False
        if not new_docs:
            self.documents = []
            self.vocab = {}
            self.idf = {}
            self.vectors = None
            for fname in ["index.json", "vectors.npy", "metadata.pkl"]:
                self._delete_oss_object(self._oss_key(fname))
            for f in [self.index_file, self.vectors_file, self.metadata_file]:
                if os.path.exists(f):
                    os.remove(f)
            self._cache_valid = False
            return True
        # 重建索引
        texts = [doc["text"] for doc in new_docs]
        self._build_vocab_and_idf(texts)
        vectors = [self._text_to_vector(t) for t in texts]
        self.vectors = np.array(vectors)
        self.documents = new_docs
        with open(self.index_file, 'w', encoding='utf-8') as f:
            json.dump({
                "documents": self.documents,
                "vocab": self.vocab,
                "idf": self.idf
            }, f, ensure_ascii=False)
        np.save(self.vectors_file, self.vectors)
        self._upload_to_oss(self.index_file, self._oss_key("index.json"))
        self._upload_to_oss(self.vectors_file, self._oss_key("vectors.npy"))
        self._cache_valid = True
        return True

    def clear_index(self):
        self.documents = []
        self.vocab = {}
        self.idf = {}
        self.vectors = None
        for fname in ["index.json", "vectors.npy", "metadata.pkl"]:
            self._delete_oss_object(self._oss_key(fname))
        for f in [self.index_file, self.vectors_file, self.metadata_file]:
            if os.path.exists(f):
                os.remove(f)
        self._cache_valid = False

    def update_metadata(self, files_info: List[Dict]):
        """单独更新 metadata.pkl 到 OSS"""
        if not files_info:
            self._delete_oss_object(self._oss_key("metadata.pkl"))
            return
        with open(self.metadata_file, 'wb') as f:
            pickle.dump({"files_info": files_info, "updated": datetime.now().isoformat()}, f)
        self._upload_to_oss(self.metadata_file, self._oss_key("metadata.pkl"))


# ============================== RDM 单号处理 ============================
class RDMService:
    @staticmethod
    def fetch_rdm_content(rdm_code: str) -> Dict:
        """获取RDM单号对应的描述内容"""
        try:
            url = f"{AppConfig.RDM_BASE_URL}{rdm_code}"
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
            }
            response = requests.get(url, headers=headers, timeout=30)

            if response.status_code == 200:
                soup = BeautifulSoup(response.text, 'html.parser')
                # 查找描述模块
                description_div = soup.find('div', id='description-val')
                if description_div:
                    # 提取用户内容块
                    user_content = description_div.find('div', class_='user-content-block')
                    if user_content:
                        content = user_content.get_text(strip=True)
                        return {"success": True, "content": content, "url": url}
                    else:
                        content = description_div.get_text(strip=True)
                        return {"success": True, "content": content, "url": url}
                else:
                    return {"success": False, "error": "未找到描述内容", "url": url}
            else:
                return {"success": False, "error": f"HTTP {response.status_code}", "url": url}
        except Exception as e:
            return {"success": False, "error": str(e), "url": url}

    @staticmethod
    def extract_rdm_codes(text: str) -> List[str]:
        """从文本中提取RDM单号：字母开头 + 字母/数字/下划线 + - + 数字"""
        pattern = r'[A-Za-z][A-Za-z0-9_]*-\d+'
        return re.findall(pattern, text)

#===================================测试用例解析与导出================
class TestCaseService:
    @staticmethod
    def parse(content: str, rdm_codes: List[str]) -> List[Dict]:
        """
        解析管道符分隔的测试用例格式
        格式：用例ID|优先级|用例名称|前置条件|测试步骤|预期结果
        步骤和预期结果使用分号(;)或中文分号(；)分隔
        """
        cases = []
        lines = content.strip().splitlines()
        rdm_idx = 0

        for line in lines:
            line = line.strip()
            if not line:
                continue
            # 按 | 分割成6个字段
            parts = line.split('|')
            if len(parts) < 6:
                # 尝试兼容可能包含管道符的内容
                if len(parts) == 5:
                    # 缺少优先级，补默认值
                    parts.insert(1, "Medium")
                elif len(parts) > 6:
                    # 合并多余的字段到预期结果
                    parts[5] = '|'.join(parts[5:])
                    parts = parts[:6]
                else:
                    continue
            # 字段一一对应
            case_id = parts[0].strip()
            priority = parts[1].strip()
            case_name = parts[2].strip()
            pre_condition = parts[3].strip() if parts[3].strip() else "无"
            steps_text = parts[4].strip()
            expect_text = parts[5].strip()
            # 跳过无效用例
            if not case_id or not case_name:
                continue
            # 处理步骤和预期结果（保留原始分隔符，不转义）
            test_steps = TestCaseService._format_steps(steps_text)
            expected_result = TestCaseService._format_expected(expect_text)
            # 组装成字典
            case = {
                "RDM单号": "",
                "用例ID": case_id,
                "优先级": priority,
                "用例名称": case_name,
                "前置条件": pre_condition,
                "测试步骤": test_steps,
                "预期结果": expected_result
            }

            # 自动填充 RDM 单号
            if rdm_codes:
                case["RDM单号"] = rdm_codes[rdm_idx % len(rdm_codes)]
                rdm_idx += 1

            cases.append(case)

        return cases

    @staticmethod
    def _format_steps(steps_text: str) -> str:
        """
        格式化测试步骤
        输入：1.打开页面->点击按钮；2.填写表单；3.提交
        输出：1. 打开页面->点击按钮；2. 填写表单；3. 提交
        注意：不显示转义字符，保持原始分号分隔
        """
        if not steps_text:
            return ""

        # 统一分隔符：将中文分号替换为英文分号
        steps_text = steps_text.replace('；', ';')

        # 按分号分割步骤
        steps = [s.strip() for s in steps_text.split(';') if s.strip()]

        # 格式化输出：保持分号分隔，不换行
        formatted_steps = []
        for i, step in enumerate(steps, 1):
            # 移除步骤中已有的编号（如"1."），避免重复
            step = re.sub(r'^\d+\.\s*', '', step)
            formatted_steps.append(f"{i}.{step}")
        # 使用分号连接，不显示换行符
        return '；'.join(formatted_steps)

    @staticmethod
    def _format_expected(expect_text: str) -> str:
        """
        格式化预期结果
        输入：1.弹出审核意见填写框；2.状态更新为待开票
        输出：1.弹出审核意见填写框；2.状态更新为待开票
        注意：不显示转义字符，保持原始分号分隔
        """
        if not expect_text:
            return ""
        # 统一分隔符：将中文分号替换为英文分号
        expect_text = expect_text.replace('；', ';')
        # 按分号分割预期结果
        expects = [e.strip() for e in expect_text.split(';') if e.strip()]
        # 格式化输出：保持分号分隔，不换行
        formatted_expects = []
        for i, exp in enumerate(expects, 1):
            # 移除已有的编号
            exp = re.sub(r'^\d+\.\s*', '', exp)
            formatted_expects.append(f"{i}.{exp}")
        # 使用分号连接，不显示换行符
        return '；'.join(formatted_expects)

class ExportService:
    @staticmethod
    def to_csv(cases: List[Dict]) -> bytes:
        """导出为标准CSV格式，不包含转义字符"""
        if not cases:
            return b''
        df = pd.DataFrame(cases)
        cols = ["RDM单号", "用例ID", "优先级", "用例名称", "前置条件", "测试步骤", "预期结果"]
        # 只保留存在的列
        existing_cols = [c for c in cols if c in df.columns]
        df = df[existing_cols]
        # 确保不包含转义字符
        for col in df.columns:
            if df[col].dtype == 'object':
                df[col] = df[col].astype(str).str.replace('\n', ' ').str.replace('\r', ' ')
        buf = io.BytesIO()
        # 修复：使用 lineterminator 而不是 line_terminator
        df.to_csv(buf, index=False, encoding='utf-8-sig', lineterminator='\n')
        return buf.getvalue()